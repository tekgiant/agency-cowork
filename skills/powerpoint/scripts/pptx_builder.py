"""
PowerPoint Presentation Builder.

Creates .pptx files from structured content using python-pptx.
Reads branding, layout, and template settings from agentconfig.json.

Usage:
    python -m scripts.pptx_builder --json slides.json --output presentation.pptx
    python -m scripts.pptx_builder --title "Title" --subtitle "Sub" --slides-json '[...]'
"""

import argparse
import json
import os
import sys
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE

# Add office-common to path
sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent))
from office_common.config import load_config, OfficeConfig


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert '#RRGGBB' to RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class PptxBuilder:
    """Build PowerPoint presentations from structured content."""

    def __init__(self, config: Optional[OfficeConfig] = None, template: Optional[str] = None):
        self.config = config or load_config()
        self.brand = self.config.branding
        self.pptx_cfg = self.config.powerpoint

        template_path = template or self.pptx_cfg.default_template
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
            # Set slide dimensions (widescreen 16:9)
            self.prs.slide_width = Inches(self.pptx_cfg.slide_width_inches)
            self.prs.slide_height = Inches(self.pptx_cfg.slide_height_inches)

    def _apply_font(self, run, size_pt: int, bold: bool = False, color: Optional[str] = None,
                     font_name: Optional[str] = None):
        """Apply font formatting to a text run."""
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.name = font_name or self.brand.font_body
        if color:
            run.font.color.rgb = hex_to_rgb(color)

    def _add_footer(self, slide):
        """Add footer text and optional slide number to a slide."""
        if not self.pptx_cfg.footer_text and not self.pptx_cfg.slide_numbers:
            return

        from pptx.util import Inches, Pt
        footer_text = self.pptx_cfg.footer_text or ""
        if self.pptx_cfg.slide_numbers:
            slide_num = len(self.prs.slides)
            footer_text = f"{footer_text}    |    Slide {slide_num}" if footer_text else f"Slide {slide_num}"

        left = Inches(0.5)
        top = self.prs.slide_height - Inches(0.5)
        width = self.prs.slide_width - Inches(1.0)
        height = Inches(0.3)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = footer_text
        self._apply_font(run, 8, color="#888888")

    def add_title_slide(self, title: str, subtitle: str = "", notes: str = ""):
        """Add a title slide with large centered title and subtitle."""
        slide_layout = self.prs.slide_layouts[0]  # Title Slide layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        if slide.placeholders[0]:
            title_shape = slide.placeholders[0]
            title_shape.text = title
            for para in title_shape.text_frame.paragraphs:
                for run in para.runs:
                    self._apply_font(run, self.pptx_cfg.title_font_size_pt + 8, bold=True,
                                     color=self.brand.dark_color, font_name=self.brand.font_heading)

        # Subtitle
        if len(slide.placeholders) > 1 and subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            for para in subtitle_shape.text_frame.paragraphs:
                for run in para.runs:
                    self._apply_font(run, self.pptx_cfg.body_font_size_pt + 2, color="#666666")

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self._add_footer(slide)
        return slide

    def add_section_header(self, title: str, subtitle: str = "", notes: str = ""):
        """Add a section divider slide."""
        slide_layout = self.prs.slide_layouts[2]  # Section Header
        slide = self.prs.slides.add_slide(slide_layout)

        if slide.placeholders[0]:
            shape = slide.placeholders[0]
            shape.text = title
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    self._apply_font(run, self.pptx_cfg.title_font_size_pt, bold=True,
                                     color=self.brand.primary_color, font_name=self.brand.font_heading)

        if len(slide.placeholders) > 1 and subtitle:
            shape = slide.placeholders[1]
            shape.text = subtitle

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self._add_footer(slide)
        return slide

    def add_content_slide(self, title: str, bullets: list[str], notes: str = "",
                          two_column: bool = False):
        """Add a content slide with title and bullet points."""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        if slide.placeholders[0]:
            shape = slide.placeholders[0]
            shape.text = title
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    self._apply_font(run, self.pptx_cfg.title_font_size_pt, bold=True,
                                     color=self.brand.dark_color, font_name=self.brand.font_heading)

        # Content area
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()

            if two_column and len(bullets) > 3:
                mid = len(bullets) // 2
                left_bullets = bullets[:mid]
                right_bullets = bullets[mid:]
                # Use single column for now (two-column requires custom shapes)
                all_bullets = left_bullets + right_bullets
            else:
                all_bullets = bullets

            for i, bullet in enumerate(all_bullets[:self.pptx_cfg.max_bullets_per_slide]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.level = 0
                run = p.add_run()
                run.text = bullet
                self._apply_font(run, self.pptx_cfg.body_font_size_pt)

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self._add_footer(slide)
        return slide

    def add_table_slide(self, title: str, headers: list[str], rows: list[list[str]],
                        notes: str = ""):
        """Add a slide with a formatted table."""
        slide_layout = self.prs.slide_layouts[5]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        left = Inches(0.5)
        top = Inches(0.3)
        width = self.prs.slide_width - Inches(1.0)
        height = Inches(0.8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        self._apply_font(run, self.pptx_cfg.title_font_size_pt, bold=True,
                         color=self.brand.dark_color, font_name=self.brand.font_heading)

        # Table
        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)
        tbl_left = Inches(0.5)
        tbl_top = Inches(1.3)
        tbl_width = self.prs.slide_width - Inches(1.0)
        tbl_height = Inches(0.4) * num_rows

        table_shape = slide.shapes.add_table(num_rows, num_cols, tbl_left, tbl_top,
                                              tbl_width, tbl_height)
        table = table_shape.table

        # Header row
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(self.brand.primary_color)
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    self._apply_font(run, 12, bold=True, color="#FFFFFF",
                                     font_name=self.brand.font_heading)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Data rows
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(val)
                # Alternate row shading
                if i % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = hex_to_rgb("#F3F2F1")
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        self._apply_font(run, 11)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self._add_footer(slide)
        return slide

    def add_comparison_slide(self, title: str, left_title: str, left_items: list[str],
                              right_title: str, right_items: list[str], notes: str = ""):
        """Add a two-column comparison slide."""
        slide_layout = self.prs.slide_layouts[5]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3),
                                          self.prs.slide_width - Inches(1.0), Inches(0.8))
        tf = txBox.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = title
        self._apply_font(run, self.pptx_cfg.title_font_size_pt, bold=True,
                         color=self.brand.dark_color, font_name=self.brand.font_heading)

        col_width = (self.prs.slide_width - Inches(1.5)) // 2
        col_top = Inches(1.5)
        col_height = self.prs.slide_height - Inches(2.5)

        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), col_top, col_width, col_height)
        ltf = left_box.text_frame
        ltf.word_wrap = True
        lp = ltf.paragraphs[0]
        lr = lp.add_run()
        lr.text = left_title
        self._apply_font(lr, self.pptx_cfg.body_font_size_pt + 2, bold=True,
                         color=self.brand.primary_color, font_name=self.brand.font_heading)

        for item in left_items:
            p = ltf.add_paragraph()
            p.space_before = Pt(6)
            r = p.add_run()
            r.text = f"• {item}"
            self._apply_font(r, self.pptx_cfg.body_font_size_pt - 2)

        # Right column
        right_left = Inches(0.5) + col_width + Inches(0.5)
        right_box = slide.shapes.add_textbox(right_left, col_top, col_width, col_height)
        rtf = right_box.text_frame
        rtf.word_wrap = True
        rp = rtf.paragraphs[0]
        rr = rp.add_run()
        rr.text = right_title
        self._apply_font(rr, self.pptx_cfg.body_font_size_pt + 2, bold=True,
                         color=self.brand.accent_color, font_name=self.brand.font_heading)

        for item in right_items:
            p = rtf.add_paragraph()
            p.space_before = Pt(6)
            r = p.add_run()
            r.text = f"• {item}"
            self._apply_font(r, self.pptx_cfg.body_font_size_pt - 2)

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self._add_footer(slide)
        return slide

    def add_image_slide(self, title: str, image_path: str, caption: str = "",
                        notes: str = ""):
        """Add a slide with a centered image and optional caption."""
        slide_layout = self.prs.slide_layouts[5]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3),
                                          self.prs.slide_width - Inches(1.0), Inches(0.8))
        tf = txBox.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = title
        self._apply_font(run, self.pptx_cfg.title_font_size_pt, bold=True,
                         color=self.brand.dark_color, font_name=self.brand.font_heading)

        # Image (centered, max 80% width)
        if os.path.exists(image_path):
            max_w = self.prs.slide_width * 0.8
            max_h = self.prs.slide_height * 0.55
            img = slide.shapes.add_picture(image_path, Inches(0), Inches(1.3))

            # Scale to fit
            ratio = min(max_w / img.width, max_h / img.height)
            if ratio < 1:
                img.width = int(img.width * ratio)
                img.height = int(img.height * ratio)

            # Center
            img.left = int((self.prs.slide_width - img.width) / 2)

        if caption:
            cap_box = slide.shapes.add_textbox(
                Inches(0.5), self.prs.slide_height - Inches(1.2),
                self.prs.slide_width - Inches(1.0), Inches(0.4))
            p = cap_box.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = caption
            self._apply_font(run, 10, color="#888888")

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        self._add_footer(slide)
        return slide

    def add_blank_slide(self, notes: str = ""):
        """Add a blank slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        self._add_footer(slide)
        return slide

    def build_from_json(self, slides_data: list[dict]):
        """Build slides from a list of slide definitions.

        Each slide dict has:
            type: "title" | "section" | "content" | "table" | "comparison" | "image" | "blank"
            ...plus type-specific fields (title, bullets, headers, rows, etc.)
        """
        for sd in slides_data:
            slide_type = sd.get("type", "content")
            notes = sd.get("notes", "")

            if slide_type == "title":
                self.add_title_slide(sd.get("title", ""), sd.get("subtitle", ""), notes)
            elif slide_type == "section":
                self.add_section_header(sd.get("title", ""), sd.get("subtitle", ""), notes)
            elif slide_type == "content":
                self.add_content_slide(sd.get("title", ""), sd.get("bullets", []),
                                       notes, sd.get("two_column", False))
            elif slide_type == "table":
                self.add_table_slide(sd.get("title", ""), sd.get("headers", []),
                                     sd.get("rows", []), notes)
            elif slide_type == "comparison":
                self.add_comparison_slide(
                    sd.get("title", ""),
                    sd.get("left_title", "Option A"), sd.get("left_items", []),
                    sd.get("right_title", "Option B"), sd.get("right_items", []),
                    notes)
            elif slide_type == "image":
                self.add_image_slide(sd.get("title", ""), sd.get("image_path", ""),
                                     sd.get("caption", ""), notes)
            elif slide_type == "blank":
                self.add_blank_slide(notes)

    def save(self, output_path: Optional[str] = None) -> str:
        """Save the presentation and return the path."""
        if output_path is None:
            os.makedirs(self.config.output_dir, exist_ok=True)
            output_path = os.path.join(self.config.output_dir, "presentation.pptx")
        else:
            os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

        self.prs.save(output_path)
        return output_path


def main():
    parser = argparse.ArgumentParser(description="Build PowerPoint presentations")
    parser.add_argument("--json", help="Path to JSON file with slide definitions")
    parser.add_argument("--slides-json", help="Inline JSON string with slide definitions")
    parser.add_argument("--title", help="Quick: create a single title slide")
    parser.add_argument("--subtitle", default="", help="Subtitle for quick title slide")
    parser.add_argument("--output", "-o", help="Output .pptx file path")
    parser.add_argument("--template", help="Path to template .pptx file")
    parser.add_argument("--config", help="Path to agentconfig.json (default: project root)")
    args = parser.parse_args()

    config = load_config(args.config)
    builder = PptxBuilder(config, template=args.template)

    if args.json:
        with open(args.json, "r", encoding="utf-8") as f:
            slides = json.load(f)
        builder.build_from_json(slides)
    elif args.slides_json:
        slides = json.loads(args.slides_json)
        builder.build_from_json(slides)
    elif args.title:
        builder.add_title_slide(args.title, args.subtitle)
    else:
        # Read from stdin
        data = json.load(sys.stdin)
        builder.build_from_json(data)

    path = builder.save(args.output)
    print(json.dumps({"status": "ok", "path": os.path.abspath(path),
                       "slides": len(builder.prs.slides)}))


if __name__ == "__main__":
    main()
