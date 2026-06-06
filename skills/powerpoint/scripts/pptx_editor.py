"""
PowerPoint Presentation Editor.

Primary tool for inspecting and modifying existing .pptx files.
Workflow: download from SharePoint → inspect → edit → upload back.

Usage:
    python -m scripts.pptx_editor --input deck.pptx --action inspect
    python -m scripts.pptx_editor --input deck.pptx --action inspect --slide 2
    python -m scripts.pptx_editor --input deck.pptx --action replace-text --find "OLD" --replace "NEW"
    python -m scripts.pptx_editor --input deck.pptx --action update-text --slide 2 --shape 1 --text "New content"
    python -m scripts.pptx_editor --input deck.pptx --action update-table --slide 3 --shape 0 --row 1 --col 2 --text "Updated"
    python -m scripts.pptx_editor --input deck.pptx --action delete-slide --index 3
    python -m scripts.pptx_editor --input deck.pptx --action move-slide --from-index 4 --to-index 1
    python -m scripts.pptx_editor --input deck.pptx --action duplicate-slide --index 2
    python -m scripts.pptx_editor --input deck.pptx --action add-slide --slides-json '[...]'
    python -m scripts.pptx_editor --input deck.pptx --action update-notes --slide 0 --text "Speaker notes"
    python -m scripts.pptx_editor --input deck.pptx --action batch --ops-json '[...]'
"""

import argparse
import copy
import json
import os
import sys
from pathlib import Path
from typing import Optional
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu

sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent))
from office_common.config import load_config


# ─── Inspect ──────────────────────────────────────────────────────────────────

def inspect_presentation(prs: Presentation, slide_index: Optional[int] = None) -> dict:
    """Comprehensive inspection of a presentation or a single slide.

    Returns structure with enough detail for the agent to plan edits.
    """
    info = {
        "slide_count": len(prs.slides),
        "width_inches": round(prs.slide_width / 914400, 2),
        "height_inches": round(prs.slide_height / 914400, 2),
    }

    slides_to_inspect = []
    if slide_index is not None:
        if 0 <= slide_index < len(prs.slides):
            slides_to_inspect = [(slide_index, prs.slides[slide_index])]
        else:
            info["error"] = f"Slide index {slide_index} out of range (0-{len(prs.slides)-1})"
            return info
    else:
        slides_to_inspect = list(enumerate(prs.slides))

    info["slides"] = []
    for idx, slide in slides_to_inspect:
        slide_info = _inspect_slide(idx, slide)
        info["slides"].append(slide_info)

    return info


def _inspect_slide(idx: int, slide) -> dict:
    """Inspect a single slide — all shapes, text, tables, images."""
    layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"

    # Speaker notes
    notes = ""
    try:
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
    except Exception:
        pass

    shapes = []
    for si, shape in enumerate(slide.shapes):
        shape_info = {
            "shape_index": si,
            "name": shape.name,
            "type": _shape_type(shape),
            "left": round(shape.left / 914400, 2) if shape.left else None,
            "top": round(shape.top / 914400, 2) if shape.top else None,
            "width": round(shape.width / 914400, 2) if shape.width else None,
            "height": round(shape.height / 914400, 2) if shape.height else None,
        }

        # Text content
        if shape.has_text_frame:
            paragraphs = []
            for pi, para in enumerate(shape.text_frame.paragraphs):
                p_info = {
                    "para_index": pi,
                    "text": para.text,
                    "level": para.level,
                }
                if para.runs:
                    p_info["runs"] = [{"text": r.text, "bold": r.font.bold,
                                        "size_pt": round(r.font.size / 12700, 1) if r.font.size else None}
                                       for r in para.runs]
                paragraphs.append(p_info)
            shape_info["paragraphs"] = paragraphs

        # Table content
        if shape.has_table:
            table = shape.table
            table_data = {
                "rows": len(table.rows),
                "cols": len(table.columns),
                "cells": [],
            }
            for ri, row in enumerate(table.rows):
                for ci, cell in enumerate(row.cells):
                    table_data["cells"].append({
                        "row": ri, "col": ci, "text": cell.text
                    })
            shape_info["table"] = table_data

        # Image
        if shape.shape_type == 13:  # Picture
            try:
                shape_info["image"] = {
                    "content_type": shape.image.content_type,
                    "size_bytes": len(shape.image.blob),
                }
            except Exception:
                shape_info["image"] = {"note": "embedded image"}

        # Placeholder info
        if shape.is_placeholder:
            shape_info["placeholder_idx"] = shape.placeholder_format.idx
            shape_info["placeholder_type"] = str(shape.placeholder_format.type)

        shapes.append(shape_info)

    return {
        "slide_index": idx,
        "layout": layout_name,
        "shapes": shapes,
        "notes": notes,
    }


def _shape_type(shape) -> str:
    """Return a human-readable shape type."""
    if shape.has_table:
        return "table"
    if shape.shape_type == 13:
        return "image"
    if shape.has_text_frame:
        if shape.is_placeholder:
            return "placeholder"
        return "textbox"
    if hasattr(shape, "chart"):
        return "chart"
    return str(shape.shape_type)


# ─── Edit Operations ──────────────────────────────────────────────────────────

def replace_text(prs: Presentation, find: str, replace: str,
                 slide_index: Optional[int] = None) -> int:
    """Replace text across slides. Returns count of replacements."""
    count = 0
    slides = [prs.slides[slide_index]] if slide_index is not None else prs.slides
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if find in run.text:
                            run.text = run.text.replace(find, replace)
                            count += 1
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if find in run.text:
                                    run.text = run.text.replace(find, replace)
                                    count += 1
    return count


def update_shape_text(prs: Presentation, slide_idx: int, shape_idx: int,
                      text: str, para_idx: Optional[int] = None) -> bool:
    """Update text in a specific shape. If para_idx given, updates just that paragraph."""
    slide = prs.slides[slide_idx]
    shape = slide.shapes[shape_idx]
    if not shape.has_text_frame:
        return False

    if para_idx is not None:
        para = shape.text_frame.paragraphs[para_idx]
        # Preserve formatting of first run
        if para.runs:
            para.runs[0].text = text
            for r in para.runs[1:]:
                r.text = ""
        else:
            para.text = text
    else:
        # Replace all text, preserving first paragraph's formatting
        tf = shape.text_frame
        if tf.paragraphs and tf.paragraphs[0].runs:
            first_run = tf.paragraphs[0].runs[0]
            # Clear all paragraphs after the first
            for i in range(len(tf.paragraphs) - 1, 0, -1):
                p_elem = tf.paragraphs[i]._p
                p_elem.getparent().remove(p_elem)
            first_run.text = text
            for r in tf.paragraphs[0].runs[1:]:
                r.text = ""
        else:
            tf.text = text
    return True


def update_table_cell(prs: Presentation, slide_idx: int, shape_idx: int,
                      row: int, col: int, text: str) -> bool:
    """Update a specific table cell."""
    slide = prs.slides[slide_idx]
    shape = slide.shapes[shape_idx]
    if not shape.has_table:
        return False
    cell = shape.table.cell(row, col)
    # Preserve formatting
    if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
        cell.text_frame.paragraphs[0].runs[0].text = text
        for r in cell.text_frame.paragraphs[0].runs[1:]:
            r.text = ""
    else:
        cell.text = text
    return True


def update_notes(prs: Presentation, slide_idx: int, text: str) -> bool:
    """Update speaker notes for a slide."""
    slide = prs.slides[slide_idx]
    if not slide.has_notes_slide:
        slide.notes_slide  # Creates notes slide
    slide.notes_slide.notes_text_frame.text = text
    return True


def delete_slide(prs: Presentation, index: int):
    """Delete slide at the given index."""
    rId = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[index]


def move_slide(prs: Presentation, from_idx: int, to_idx: int):
    """Move a slide from one position to another."""
    sldIdLst = prs.slides._sldIdLst
    el = sldIdLst[from_idx]
    sldIdLst.remove(el)
    if to_idx >= len(sldIdLst):
        sldIdLst.append(el)
    else:
        sldIdLst.insert(to_idx, el)


def duplicate_slide(prs: Presentation, index: int):
    """Duplicate a slide (appends copy at end)."""
    template = prs.slides[index]
    slide_layout = template.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Copy shapes by cloning XML
    for shape in template.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    # Copy notes
    try:
        if template.has_notes_slide:
            new_slide.notes_slide.notes_text_frame.text = template.notes_slide.notes_text_frame.text
    except Exception:
        pass

    return len(prs.slides) - 1


def extract_text(prs: Presentation) -> list[dict]:
    """Extract all text from each slide (quick summary)."""
    results = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        texts.append(cell.text)
        results.append({"slide_index": i, "text": "\n".join(texts)})
    return results


# ─── Batch Operations ─────────────────────────────────────────────────────────

def batch_ops(prs: Presentation, operations: list[dict]) -> list[dict]:
    """Execute a batch of edit operations. Each op has 'action' + params.

    Supported actions:
        replace-text: {find, replace, slide (optional)}
        update-text: {slide, shape, text, para (optional)}
        update-table: {slide, shape, row, col, text}
        update-notes: {slide, text}
        delete-slide: {index}
        move-slide: {from, to}
        duplicate-slide: {index}
    """
    results = []
    for op in operations:
        action = op.get("action")
        try:
            if action == "replace-text":
                count = replace_text(prs, op["find"], op["replace"], op.get("slide"))
                results.append({"action": action, "ok": True, "replacements": count})
            elif action == "update-text":
                ok = update_shape_text(prs, op["slide"], op["shape"], op["text"], op.get("para"))
                results.append({"action": action, "ok": ok})
            elif action == "update-table":
                ok = update_table_cell(prs, op["slide"], op["shape"], op["row"], op["col"], op["text"])
                results.append({"action": action, "ok": ok})
            elif action == "update-notes":
                ok = update_notes(prs, op["slide"], op["text"])
                results.append({"action": action, "ok": ok})
            elif action == "delete-slide":
                delete_slide(prs, op["index"])
                results.append({"action": action, "ok": True})
            elif action == "move-slide":
                move_slide(prs, op["from"], op["to"])
                results.append({"action": action, "ok": True})
            elif action == "duplicate-slide":
                new_idx = duplicate_slide(prs, op["index"])
                results.append({"action": action, "ok": True, "new_index": new_idx})
            else:
                results.append({"action": action, "ok": False, "error": f"Unknown action: {action}"})
        except Exception as e:
            results.append({"action": action, "ok": False, "error": str(e)})
    return results


# ─── CLI ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Inspect and edit PowerPoint presentations")
    parser.add_argument("--input", "-i", required=True, help="Input .pptx file")
    parser.add_argument("--output", "-o", help="Output file (default: overwrites input)")
    parser.add_argument("--action", required=True,
                        choices=["inspect", "extract-text", "replace-text",
                                 "update-text", "update-table", "update-notes",
                                 "delete-slide", "move-slide", "duplicate-slide",
                                 "add-slide", "batch"],
                        help="Action to perform")
    # Common args
    parser.add_argument("--slide", type=int, help="Slide index (0-based)")
    parser.add_argument("--shape", type=int, help="Shape index within slide")
    parser.add_argument("--para", type=int, help="Paragraph index within shape")
    parser.add_argument("--text", help="Text content")
    # Replace args
    parser.add_argument("--find", help="Text to find")
    parser.add_argument("--replace", help="Replacement text")
    # Table args
    parser.add_argument("--row", type=int, help="Table row index")
    parser.add_argument("--col", type=int, help="Table column index")
    # Move/delete args
    parser.add_argument("--index", type=int, help="Slide index for delete/duplicate")
    parser.add_argument("--from-index", type=int, help="Source index for move")
    parser.add_argument("--to-index", type=int, help="Destination index for move")
    # JSON args
    parser.add_argument("--slides-json", help="JSON slide definitions (for add-slide)")
    parser.add_argument("--ops-json", help="JSON batch operations")
    parser.add_argument("--config", help="Path to agentconfig.json")
    args = parser.parse_args()

    prs = Presentation(args.input)
    output_path = args.output or args.input

    if args.action == "inspect":
        result = inspect_presentation(prs, slide_index=args.slide)
        print(json.dumps(result, indent=2, default=str))
        return

    elif args.action == "extract-text":
        result = extract_text(prs)
        print(json.dumps(result, indent=2))
        return

    elif args.action == "replace-text":
        if not args.find:
            print("Error: --find required", file=sys.stderr)
            sys.exit(1)
        count = replace_text(prs, args.find, args.replace or "", args.slide)
        prs.save(output_path)
        print(json.dumps({"status": "ok", "replacements": count}))

    elif args.action == "update-text":
        if args.slide is None or args.shape is None or not args.text:
            print("Error: --slide, --shape, --text required", file=sys.stderr)
            sys.exit(1)
        ok = update_shape_text(prs, args.slide, args.shape, args.text, args.para)
        prs.save(output_path)
        print(json.dumps({"status": "ok" if ok else "failed"}))

    elif args.action == "update-table":
        if None in (args.slide, args.shape, args.row, args.col) or not args.text:
            print("Error: --slide, --shape, --row, --col, --text required", file=sys.stderr)
            sys.exit(1)
        ok = update_table_cell(prs, args.slide, args.shape, args.row, args.col, args.text)
        prs.save(output_path)
        print(json.dumps({"status": "ok" if ok else "failed"}))

    elif args.action == "update-notes":
        if args.slide is None or not args.text:
            print("Error: --slide and --text required", file=sys.stderr)
            sys.exit(1)
        ok = update_notes(prs, args.slide, args.text)
        prs.save(output_path)
        print(json.dumps({"status": "ok" if ok else "failed"}))

    elif args.action == "delete-slide":
        if args.index is None:
            print("Error: --index required", file=sys.stderr)
            sys.exit(1)
        delete_slide(prs, args.index)
        prs.save(output_path)
        print(json.dumps({"status": "ok", "remaining_slides": len(prs.slides)}))

    elif args.action == "move-slide":
        if args.from_index is None or args.to_index is None:
            print("Error: --from-index and --to-index required", file=sys.stderr)
            sys.exit(1)
        move_slide(prs, args.from_index, args.to_index)
        prs.save(output_path)
        print(json.dumps({"status": "ok"}))

    elif args.action == "duplicate-slide":
        if args.index is None:
            print("Error: --index required", file=sys.stderr)
            sys.exit(1)
        new_idx = duplicate_slide(prs, args.index)
        prs.save(output_path)
        print(json.dumps({"status": "ok", "new_index": new_idx,
                           "total_slides": len(prs.slides)}))

    elif args.action == "add-slide":
        if not args.slides_json:
            print("Error: --slides-json required", file=sys.stderr)
            sys.exit(1)
        from scripts.pptx_builder import PptxBuilder
        config = load_config(args.config)
        builder = PptxBuilder(config)
        builder.prs = prs
        slides = json.loads(args.slides_json)
        builder.build_from_json(slides)
        prs.save(output_path)
        print(json.dumps({"status": "ok", "total_slides": len(prs.slides)}))

    elif args.action == "batch":
        if not args.ops_json:
            print("Error: --ops-json required", file=sys.stderr)
            sys.exit(1)
        ops = json.loads(args.ops_json)
        results = batch_ops(prs, ops)
        prs.save(output_path)
        print(json.dumps({"status": "ok", "results": results,
                           "total_slides": len(prs.slides)}))


if __name__ == "__main__":
    main()
